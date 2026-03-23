#!/usr/bin/env python3
"""
PPTX → WebP slide images + manifest
Usage: python scripts/convert_slides.py
Scans slides/**/*.pptx, converts each to a folder of WebP images + slides.json
"""
import json
import subprocess
import sys
from pathlib import Path

def check_deps():
    """Check required tools are available."""
    try:
        subprocess.run(['libreoffice', '--version'], capture_output=True, check=True)
    except FileNotFoundError:
        print("ERROR: LibreOffice not found. Install with:")
        print("  Ubuntu: sudo apt-get install -y libreoffice")
        print("  macOS:  brew install --cask libreoffice")
        sys.exit(1)

    try:
        from pdf2image import convert_from_path  # noqa: F401
    except ImportError:
        print("ERROR: pdf2image not found. Install with: pip install pdf2image")
        sys.exit(1)

    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        print("WARNING: python-pptx not found. Slide notes will not be extracted.")
        print("  Install with: pip install python-pptx")


def extract_notes(pptx_path):
    """Extract speaker notes from PPTX. Returns list of strings."""
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        notes = []
        for slide in prs.slides:
            note = ""
            if slide.has_notes_slide:
                note = slide.notes_slide.notes_text_frame.text.strip()
            notes.append(note)
        return notes
    except Exception:
        return []


def convert_pptx(pptx_path):
    """Convert a single PPTX to WebP images + slides.json manifest."""
    from pdf2image import convert_from_path

    stem = pptx_path.stem
    out_dir = pptx_path.parent / stem
    out_dir.mkdir(exist_ok=True)

    print(f"  [1/4] Extracting notes from {pptx_path.name}...")
    notes = extract_notes(pptx_path)

    print(f"  [2/4] Converting to PDF via LibreOffice...")
    subprocess.run([
        'libreoffice', '--headless', '--convert-to', 'pdf',
        str(pptx_path), '--outdir', str(out_dir)
    ], check=True, capture_output=True)

    pdf_path = out_dir / f"{stem}.pdf"
    if not pdf_path.exists():
        print(f"  ERROR: PDF not generated for {pptx_path.name}")
        return None

    print(f"  [3/4] Converting PDF to images...")
    images = convert_from_path(str(pdf_path), dpi=200)
    pdf_path.unlink()

    print(f"  [4/4] Saving {len(images)} slides as WebP...")
    manifest = {
        "id": stem,
        "title": stem.replace("-", " ").replace("_", " ").title(),
        "totalSlides": len(images),
        "slides": []
    }

    for i, img in enumerate(images):
        idx = f"{i + 1:02d}"

        # Full size WebP
        webp_path = out_dir / f"slide-{idx}.webp"
        img.save(str(webp_path), 'WebP', quality=85)

        # Thumbnail (480px wide)
        thumb = img.copy()
        thumb.thumbnail((480, 270))
        thumb_path = out_dir / f"slide-{idx}-thumb.webp"
        thumb.save(str(thumb_path), 'WebP', quality=70)

        manifest["slides"].append({
            "index": i + 1,
            "image": f"slide-{idx}.webp",
            "thumbnail": f"slide-{idx}-thumb.webp",
            "notes": notes[i] if i < len(notes) else ""
        })

    # Write manifest
    manifest_path = out_dir / "slides.json"
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding='utf-8'
    )

    print(f"  ✅ {manifest['totalSlides']} slides → {out_dir}/")
    return manifest


def main():
    check_deps()

    slides_dir = Path("slides")
    if not slides_dir.exists():
        print("No slides/ directory found.")
        return

    pptx_files = [p for p in slides_dir.rglob("*.pptx") if "~$" not in p.name]
    if not pptx_files:
        print("No .pptx files found in slides/")
        return

    print(f"Found {len(pptx_files)} PPTX file(s)")
    print()

    for pptx in sorted(pptx_files):
        # Skip if already converted and PPTX hasn't changed
        out_dir = pptx.parent / pptx.stem
        manifest_path = out_dir / "slides.json"
        if manifest_path.exists():
            manifest_mtime = manifest_path.stat().st_mtime
            pptx_mtime = pptx.stat().st_mtime
            if manifest_mtime > pptx_mtime:
                print(f"⏭ Skipping {pptx.name} (already converted)")
                continue

        print(f"📄 Converting: {pptx}")
        convert_pptx(pptx)
        print()

    # Generate master index of all decks
    all_decks = []
    for manifest_path in slides_dir.rglob("slides.json"):
        deck = json.loads(manifest_path.read_text(encoding='utf-8'))
        deck["path"] = str(manifest_path.parent.relative_to(slides_dir))
        all_decks.append({
            "id": deck["id"],
            "title": deck["title"],
            "totalSlides": deck["totalSlides"],
            "path": deck["path"],
            "thumbnail": f"{deck['path']}/{deck['slides'][0]['thumbnail']}" if deck['slides'] else ""
        })

    index_path = slides_dir / "index.json"
    index_path.write_text(
        json.dumps(all_decks, ensure_ascii=False, indent=2),
        encoding='utf-8'
    )
    print(f"📋 Master index: {index_path} ({len(all_decks)} decks)")


if __name__ == "__main__":
    main()
